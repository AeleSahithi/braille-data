# scripts/extract_clean.py
import os
import sys

# Add the project root directory to Python path
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.append(project_root)

from config import PATHS, SUPPORTED_FORMATS, OCR_CONFIG
import pytesseract
import cv2
import PyPDF2
import docx
import pandas as pd
from pathlib import Path
from bs4 import BeautifulSoup
import extract_msg
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
import email
from email import policy
import csv

input_dir = "data/raw"
output_dir = "data/processed"

# Create output directory if it doesn't exist
os.makedirs(output_dir, exist_ok=True)

print(f"Processing files from {input_dir} to {output_dir}")

def clean_text(text):
    """Clean and normalize text content"""
    return text.replace("\n", " ").strip()

def process_pdf(file_path):
    """Extract text from PDF file"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error reading PDF {file_path}: {str(e)}")
    return clean_text(text)

def process_word(file_path):
    """Extract text from Word document"""
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error reading Word document {file_path}: {str(e)}")
    return clean_text(text)

def process_excel(file_path):
    """Extract text from Excel file"""
    text = ""
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            text += f"Sheet: {sheet_name}\n"
            text += df.to_string() + "\n\n"
    except Exception as e:
        print(f"Error reading Excel file {file_path}: {str(e)}")
    return clean_text(text)

def process_html(file_path):
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            text = soup.get_text(separator=' ', strip=True)
    except Exception as e:
        print(f"Error reading HTML file {file_path}: {str(e)}")
    return clean_text(text)

def process_msg(file_path):
    text = ""
    try:
        msg = extract_msg.Message(file_path)
        text = msg.body
    except Exception as e:
        print(f"Error reading MSG file {file_path}: {str(e)}")
    return clean_text(text)

def process_eml(file_path):
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            msg = email.message_from_file(f, policy=policy.default)
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        text += part.get_payload(decode=True).decode(errors='ignore')
            else:
                text = msg.get_payload(decode=True).decode(errors='ignore')
    except Exception as e:
        print(f"Error reading EML file {file_path}: {str(e)}")
    return clean_text(text)

def process_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX file {file_path}: {str(e)}")
    return clean_text(text)

def process_csv(file_path):
    text = ""
    try:
        df = pd.read_csv(file_path)
        text = df.to_string()
    except Exception as e:
        print(f"Error reading CSV file {file_path}: {str(e)}")
    return clean_text(text)

def process_rtf(file_path):
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            rtf_content = f.read()
            text = rtf_to_text(rtf_content)
    except Exception as e:
        print(f"Error reading RTF file {file_path}: {str(e)}")
    return clean_text(text)

for file in os.listdir(input_dir):
    file_path = os.path.join(input_dir, file)
    file_extension = Path(file).suffix.lower()
    
    try:
        if file_extension in ['.png', '.jpg', '.jpeg']:
            # Process image files
            print(f"Processing image: {file}")
            image = cv2.imread(file_path)
            if image is None:
                print(f"Error: Could not read image {file}")
                continue
                
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            text = pytesseract.image_to_string(gray)
            cleaned_text = clean_text(text)
            
        elif file_extension == '.pdf':
            # Process PDF files
            print(f"Processing PDF: {file}")
            cleaned_text = process_pdf(file_path)
            
        elif file_extension in ['.doc', '.docx']:
            # Process Word documents
            print(f"Processing Word document: {file}")
            cleaned_text = process_word(file_path)
            
        elif file_extension in ['.xls', '.xlsx']:
            # Process Excel files
            print(f"Processing Excel file: {file}")
            cleaned_text = process_excel(file_path)
            
        elif file_extension == '.txt':
            # Process text files
            print(f"Processing text file: {file}")
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            cleaned_text = clean_text(text)
            
        elif file_extension in ['.html', '.htm']:
            # Process HTML files
            print(f"Processing HTML file: {file}")
            cleaned_text = process_html(file_path)
            
        elif file_extension == '.msg':
            # Process MSG email files
            print(f"Processing MSG email file: {file}")
            cleaned_text = process_msg(file_path)
            
        elif file_extension == '.eml':
            # Process EML email files
            print(f"Processing EML email file: {file}")
            cleaned_text = process_eml(file_path)
            
        elif file_extension == '.pptx':
            # Process PowerPoint files
            print(f"Processing PowerPoint file: {file}")
            cleaned_text = process_pptx(file_path)
            
        elif file_extension == '.csv':
            # Process CSV files
            print(f"Processing CSV file: {file}")
            cleaned_text = process_csv(file_path)
            
        elif file_extension == '.rtf':
            # Process RTF files
            print(f"Processing RTF file: {file}")
            cleaned_text = process_rtf(file_path)
            
        else:
            print(f"Skipping {file} (not a supported file type)")
            continue
            
        # Save the extracted text
        output_file = os.path.join(output_dir, f"{file}.txt")
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(cleaned_text)
        print(f"Saved extracted text to: {output_file}")
            
    except Exception as e:
        print(f"Error processing file {file}: {str(e)}")

print("Processing complete!")
